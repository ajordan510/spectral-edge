"""
Report Generator Module for SpectralEdge.

This module provides functionality to generate PowerPoint reports containing
PSD analysis results, including plots, parameters, and summary tables.

Author: SpectralEdge Development Team
"""

import io
import tempfile
from pathlib import Path
from typing import Optional, List, Dict, Any, Tuple
from datetime import datetime

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.parts.image import Image
    from pptx.oxml.ns import qn
    from pptx.oxml.xmlchemy import OxmlElement
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

import numpy as np


class ReportGenerator:
    """
    Generate PowerPoint reports for PSD analysis results.

    This class provides methods to create professional reports containing:
    - PSD plots (captured as images)
    - Analysis parameters
    - RMS values and summary tables
    - Spectrogram images (optional)

    Attributes
    ----------
    presentation : Presentation
        The python-pptx Presentation object being built.
    title : str
        Report title.

    Examples
    --------
    >>> generator = ReportGenerator(title="Vibration Analysis Report")
    >>> generator.add_title_slide(subtitle="Flight Test 001")
    >>> generator.add_psd_plot(image_bytes, "Channel 1 PSD", {"df": 1.0, "window": "hann"})
    >>> generator.add_summary_table(channels, rms_values)
    >>> generator.save("report.pptx")
    """

    def __init__(self, title: str = "PSD Analysis Report", template_path: Optional[str] = None):
        """
        Initialize the report generator.

        Parameters
        ----------
        title : str, optional
            Report title. Default is "PSD Analysis Report".
        template_path : str, optional
            Path to a PowerPoint template file. If None, uses default blank template.
        """
        if not PPTX_AVAILABLE:
            raise ImportError(
                "python-pptx is required for report generation. "
                "Install it with: pip install python-pptx"
            )

        self.title = title

        if template_path and Path(template_path).exists():
            self.presentation = Presentation(template_path)
        else:
            self.presentation = Presentation()
            # Set default slide dimensions (16:9 widescreen)
            self.presentation.slide_width = Inches(13.333)
            self.presentation.slide_height = Inches(7.5)

        self._slide_count = 0

    def _style_slide_title(self, paragraph, size_pt: float = 28.0) -> None:
        """Apply consistent title typography across report slides."""
        paragraph.font.name = "Arial"
        paragraph.font.size = Pt(size_pt)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(0x00, 0x39, 0xA6)

    def _add_picture_fit(self, slide, image_bytes: bytes, left, top, width, height):
        """
        Add image to fit inside the target box while preserving aspect ratio.

        Prevents font and axis distortion caused by forced width/height stretching.
        """
        image = Image.from_blob(image_bytes)
        img_w_px, img_h_px = image.size
        if img_w_px <= 0 or img_h_px <= 0:
            return slide.shapes.add_picture(io.BytesIO(image_bytes), left, top, width, height)

        slot_w = int(width)
        slot_h = int(height)
        scale = min(slot_w / float(img_w_px), slot_h / float(img_h_px))
        draw_w = int(img_w_px * scale)
        draw_h = int(img_h_px * scale)
        draw_left = int(left) + (slot_w - draw_w) // 2
        draw_top = int(top) + (slot_h - draw_h) // 2

        return slide.shapes.add_picture(io.BytesIO(image_bytes), draw_left, draw_top, draw_w, draw_h)

    def _apply_medium_style_3(self, table) -> bool:
        """
        Best-effort apply PowerPoint table style "Medium Style 3".

        Returns True if style assignment succeeded, False otherwise.
        """
        # Candidate GUIDs across Office themes; first valid one wins.
        candidate_ids = [
            "{D7AC3CCA-C797-4891-BE02-D94E43425B78}",
            "{5C22544A-7EE6-4342-B048-85BDC9FD1C3A}",
        ]
        try:
            tbl_pr = table._tbl.tblPr
            if tbl_pr is None:
                return False
            style_node = tbl_pr.find(qn("a:tableStyleId"))
            if style_node is None:
                style_node = OxmlElement("a:tableStyleId")
                tbl_pr.append(style_node)
            style_node.text = candidate_ids[0]
            return True
        except Exception:
            return False

    def add_title_slide(
        self,
        subtitle: str = "",
        author: str = "",
        date: Optional[str] = None
    ) -> None:
        """
        Add a title slide to the report.

        Parameters
        ----------
        subtitle : str, optional
            Subtitle text (e.g., flight name, test ID).
        author : str, optional
            Author name.
        date : str, optional
            Date string. If None, uses current date.
        """
        # Use blank layout and add shapes manually for more control
        slide_layout = self.presentation.slide_layouts[6]  # Blank layout
        slide = self.presentation.slides.add_slide(slide_layout)

        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = self.title
        self._style_slide_title(title_para, size_pt=28.0)
        title_para.alignment = PP_ALIGN.CENTER

        # Add subtitle if provided
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(4.0), Inches(12.333), Inches(0.75)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.text = subtitle
            subtitle_para.font.size = Pt(28)
            subtitle_para.font.color.rgb = RGBColor(0x4a, 0x55, 0x68)  # Gray
            subtitle_para.alignment = PP_ALIGN.CENTER

        # Add date
        if date is None:
            date = datetime.now().strftime("%B %d, %Y")

        date_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(5.5), Inches(12.333), Inches(0.5)
        )
        date_frame = date_box.text_frame
        date_para = date_frame.paragraphs[0]
        date_para.text = date
        date_para.font.size = Pt(18)
        date_para.font.color.rgb = RGBColor(0x6b, 0x72, 0x80)  # Light gray
        date_para.alignment = PP_ALIGN.CENTER

        # Add author if provided
        if author:
            author_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(6.0), Inches(12.333), Inches(0.5)
            )
            author_frame = author_box.text_frame
            author_para = author_frame.paragraphs[0]
            author_para.text = author
            author_para.font.size = Pt(16)
            author_para.font.color.rgb = RGBColor(0x6b, 0x72, 0x80)
            author_para.alignment = PP_ALIGN.CENTER

        self._slide_count += 1

    def add_psd_plot(
        self,
        image_bytes: bytes,
        title: str = "PSD Analysis",
        parameters: Optional[Dict[str, Any]] = None,
        rms_values: Optional[Dict[str, float]] = None,
        units: str = ""
    ) -> None:
        """
        Add a slide with a PSD plot image and optional parameters.

        Parameters
        ----------
        image_bytes : bytes
            PNG image data of the PSD plot.
        title : str, optional
            Slide title.
        parameters : dict, optional
            Dictionary of analysis parameters to display.
            Example: {"df": 1.0, "window": "hann", "method": "Maximax"}
        rms_values : dict, optional
            Dictionary of channel names to RMS values.
            Example: {"Channel 1": 2.45, "Channel 2": 1.89}
        units : str, optional
            Units for RMS values (e.g., "g", "m/s²").
        """
        slide_layout = self.presentation.slide_layouts[6]  # Blank layout
        slide = self.presentation.slides.add_slide(slide_layout)

        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        self._style_slide_title(title_para)

        # Calculate image position (centered, with room for parameters)
        if parameters or rms_values:
            # Image on left, parameters on right
            img_left = Inches(0.3)
            img_top = Inches(1.0)
            img_width = Inches(9.0)
            img_height = Inches(5.5)
        else:
            # Image centered and larger
            img_left = Inches(1.0)
            img_top = Inches(1.0)
            img_width = Inches(11.0)
            img_height = Inches(6.0)

        self._add_picture_fit(slide, image_bytes, img_left, img_top, img_width, img_height)

        # Add parameters box if provided
        if parameters or rms_values:
            param_box = slide.shapes.add_textbox(
                Inches(9.5), Inches(1.0), Inches(3.5), Inches(5.5)
            )
            param_frame = param_box.text_frame
            param_frame.word_wrap = True

            # Add parameters header
            if parameters:
                para = param_frame.paragraphs[0]
                para.text = "Parameters"
                para.font.size = Pt(14)
                para.font.bold = True
                para.font.color.rgb = RGBColor(0x1a, 0x1f, 0x2e)

                # Add each parameter
                for key, value in parameters.items():
                    para = param_frame.add_paragraph()
                    para.text = f"{key}: {value}"
                    para.font.size = Pt(11)
                    para.font.color.rgb = RGBColor(0x4a, 0x55, 0x68)

                # Add spacing
                para = param_frame.add_paragraph()
                para.text = ""

            # Add RMS values if provided
            if rms_values:
                para = param_frame.add_paragraph() if parameters else param_frame.paragraphs[0]
                para.text = "RMS Values"
                para.font.size = Pt(14)
                para.font.bold = True
                para.font.color.rgb = RGBColor(0x1a, 0x1f, 0x2e)

                for channel, rms in rms_values.items():
                    para = param_frame.add_paragraph()
                    if units:
                        para.text = f"{channel}: {rms:.4f} {units}"
                    else:
                        para.text = f"{channel}: {rms:.4f}"
                    para.font.size = Pt(11)
                    para.font.color.rgb = RGBColor(0x4a, 0x55, 0x68)

        self._slide_count += 1

    def add_single_plot_slide(
        self,
        image_bytes: bytes,
        title: str
    ) -> None:
        """Add a slide with a single full-width plot image."""
        slide_layout = self.presentation.slide_layouts[6]
        slide = self.presentation.slides.add_slide(slide_layout)

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        self._style_slide_title(title_para)

        # Image
        self._add_picture_fit(slide, image_bytes, Inches(0.5), Inches(1.0), Inches(12.333), Inches(6.0))

        self._slide_count += 1

    def add_two_plot_slide(
        self,
        title: str,
        left_image: bytes,
        right_image: bytes
    ) -> None:
        """Add a slide with two side-by-side plots."""
        slide_layout = self.presentation.slide_layouts[6]
        slide = self.presentation.slides.add_slide(slide_layout)

        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        self._style_slide_title(title_para)

        self._add_picture_fit(slide, left_image, Inches(0.5), Inches(1.1), Inches(6.1), Inches(5.9))
        self._add_picture_fit(slide, right_image, Inches(6.9), Inches(1.1), Inches(6.1), Inches(5.9))

        self._slide_count += 1

    def add_three_plot_slide(
        self,
        title: str,
        top_image: bytes,
        bottom_left_image: bytes,
        bottom_right_image: bytes
    ) -> None:
        """Add a slide with a top plot and two bottom plots."""
        slide_layout = self.presentation.slide_layouts[6]
        slide = self.presentation.slides.add_slide(slide_layout)

        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        self._style_slide_title(title_para)

        # Geometry tuned for Time+PSD+Spectrogram layout:
        # bottom panels are 30% taller; top panel is reduced to maintain slide fit.
        self._add_picture_fit(slide, top_image, Inches(0.5), Inches(0.95), Inches(12.333), Inches(1.92))
        self._add_picture_fit(slide, bottom_left_image, Inches(0.5), Inches(3.05), Inches(6.1), Inches(4.03))
        self._add_picture_fit(slide, bottom_right_image, Inches(6.9), Inches(3.05), Inches(6.1), Inches(4.03))

        self._slide_count += 1

    def add_statistics_slide(
        self,
        title: str,
        pdf_image: bytes,
        stats_image: bytes,
        summary_text: str
    ) -> None:
        """Add a slide with PDF and running stats plots plus summary text."""
        slide_layout = self.presentation.slide_layouts[6]
        slide = self.presentation.slides.add_slide(slide_layout)

        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        self._style_slide_title(title_para)

        self._add_picture_fit(slide, pdf_image, Inches(0.5), Inches(1.1), Inches(6.1), Inches(4.0))
        self._add_picture_fit(slide, stats_image, Inches(6.9), Inches(1.1), Inches(6.1), Inches(4.0))

        # Summary text
        summary_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(5.3), Inches(12.333), Inches(1.7)
        )
        summary_frame = summary_box.text_frame
        summary_frame.word_wrap = True
        para = summary_frame.paragraphs[0]
        para.text = summary_text
        para.font.size = Pt(12)
        para.font.color.rgb = RGBColor(0x4a, 0x55, 0x68)

        self._slide_count += 1

    def add_rms_table_slide(
        self,
        title: str,
        headers: List[str],
        rows: List[List[str]]
    ) -> None:
        """Add a slide with a single RMS summary table."""
        slide_layout = self.presentation.slide_layouts[6]
        slide = self.presentation.slides.add_slide(slide_layout)

        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        self._style_slide_title(title_para)

        num_cols = max(1, len(headers))
        num_rows = len(rows) + 1
        table_width = Inches(12.2)
        table_left = Inches(0.5)
        table_top = Inches(1.1)
        table_height = Inches(min(0.4 * num_rows, 5.9))

        table = slide.shapes.add_table(
            num_rows, num_cols, table_left, table_top, table_width, table_height
        ).table

        for col, header in enumerate(headers):
            cell = table.cell(0, col)
            cell.text = header
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(11)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        for row_idx, row_values in enumerate(rows, start=1):
            for col_idx in range(num_cols):
                value = row_values[col_idx] if col_idx < len(row_values) else ""
                cell = table.cell(row_idx, col_idx)
                cell.text = value
                cell.text_frame.paragraphs[0].font.size = Pt(10)
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        self._slide_count += 1

    def add_statistics_dashboard_slide(
        self,
        title: str,
        pdf_image: bytes,
        mean_image: bytes,
        std_image: bytes,
        skew_image: bytes,
        kurt_image: bytes,
        summary_rows: List[Tuple[str, str]]
    ) -> None:
        """Add a dashboard-style statistics slide."""
        slide_layout = self.presentation.slide_layouts[6]
        slide = self.presentation.slides.add_slide(slide_layout)

        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        self._style_slide_title(title_para)

        # Expanded running-stats area on left, tighter PDF/table panel on right.
        left_x = Inches(0.35)
        left_w = Inches(8.35)
        top_y = Inches(0.95)
        plot_h = Inches(1.55)
        gap = Inches(0.08)

        plots = [mean_image, std_image, skew_image, kurt_image]
        for idx, img in enumerate(plots):
            y = top_y + idx * (plot_h + gap)
            self._add_picture_fit(slide, img, left_x, y, left_w, plot_h)

        right_x = Inches(8.95)
        pdf_size = Inches(3.38)
        self._add_picture_fit(slide, pdf_image, right_x, top_y, pdf_size, pdf_size)

        table_top = top_y + pdf_size + Inches(0.14)
        table_height = Inches(2.35)
        table_width = Inches(2.72)
        table_left = Inches(9.28)

        num_rows = len(summary_rows) + 1
        table = slide.shapes.add_table(
            num_rows, 2, table_left, table_top, table_width, table_height
        ).table
        # Force requested column width.
        table.columns[0].width = Inches(1.36)
        table.columns[1].width = Inches(1.36)
        self._apply_medium_style_3(table)

        headers = ["Metric", "Value"]
        for col, header in enumerate(headers):
            cell = table.cell(0, col)
            cell.text = header
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(10)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        for row_idx, (metric, value) in enumerate(summary_rows, start=1):
            table.cell(row_idx, 0).text = metric
            table.cell(row_idx, 1).text = value
            table.cell(row_idx, 0).text_frame.paragraphs[0].font.size = Pt(9)
            table.cell(row_idx, 1).text_frame.paragraphs[0].font.size = Pt(9)
            table.cell(row_idx, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        self._slide_count += 1

    def add_bulleted_sections_slide(
        self,
        title: str,
        sections: List[Tuple[str, List[str]]]
    ) -> None:
        """Add a slide with section headers and bullet lists."""
        slide_layout = self.presentation.slide_layouts[6]
        slide = self.presentation.slides.add_slide(slide_layout)

        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        self._style_slide_title(title_para)

        content_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(1.2), Inches(12.0), Inches(5.8)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True

        first = True
        for section_title, bullets in sections:
            para = content_frame.paragraphs[0] if first else content_frame.add_paragraph()
            first = False
            para.text = section_title
            para.font.size = Pt(16)
            para.font.bold = True
            para.font.color.rgb = RGBColor(0x1a, 0x1f, 0x2e)
            para.space_after = Pt(4)

            for bullet in bullets:
                bpara = content_frame.add_paragraph()
                bpara.text = f"- {bullet}"
                bpara.font.size = Pt(13)
                bpara.font.color.rgb = RGBColor(0x4a, 0x55, 0x68)
                bpara.level = 1

        self._slide_count += 1

    def add_parameters_slide(self, title: str, content: str) -> None:
        """Add a parameters slide."""
        self.add_text_slide(title=title, content=content)

    def add_spectrogram(
        self,
        image_bytes: bytes,
        title: str = "Spectrogram",
        channel_name: str = "",
        parameters: Optional[Dict[str, Any]] = None
    ) -> None:
        """
        Add a slide with a spectrogram image.

        Parameters
        ----------
        image_bytes : bytes
            PNG image data of the spectrogram.
        title : str, optional
            Slide title.
        channel_name : str, optional
            Name of the channel being displayed.
        parameters : dict, optional
            Dictionary of spectrogram parameters.
        """
        slide_layout = self.presentation.slide_layouts[6]  # Blank layout
        slide = self.presentation.slides.add_slide(slide_layout)

        # Add title
        full_title = f"{title}: {channel_name}" if channel_name else title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = full_title
        self._style_slide_title(title_para)

        # Add the spectrogram image (full width)
        self._add_picture_fit(slide, image_bytes, Inches(0.5), Inches(1.0), Inches(12.333), Inches(6.0))

        self._slide_count += 1

    def add_summary_table(
        self,
        channels: List[str],
        rms_values: Dict[str, float],
        units: str = "",
        additional_metrics: Optional[Dict[str, Dict[str, float]]] = None,
        title: str = "Analysis Summary"
    ) -> None:
        """
        Add a summary table slide with RMS values for all channels.

        Parameters
        ----------
        channels : list of str
            List of channel names.
        rms_values : dict
            Dictionary mapping channel names to RMS values.
        units : str, optional
            Units for RMS values.
        additional_metrics : dict, optional
            Additional metrics to include. Format: {"Metric Name": {"Channel": value}}
        title : str, optional
            Slide title.
        """
        slide_layout = self.presentation.slide_layouts[6]  # Blank layout
        slide = self.presentation.slides.add_slide(slide_layout)

        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        self._style_slide_title(title_para)

        # Determine table dimensions
        num_cols = 2  # Channel, RMS
        if additional_metrics:
            num_cols += len(additional_metrics)
        num_rows = len(channels) + 1  # Header + data rows

        # Create table
        table_width = Inches(10)
        table_height = Inches(0.4 * num_rows)
        table_left = Inches(1.5)
        table_top = Inches(1.5)

        table = slide.shapes.add_table(
            num_rows, num_cols, table_left, table_top, table_width, table_height
        ).table

        # Style the table
        # Header row
        headers = ["Channel", f"RMS ({units})" if units else "RMS"]
        if additional_metrics:
            headers.extend(additional_metrics.keys())

        for col, header in enumerate(headers):
            cell = table.cell(0, col)
            cell.text = header
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Data rows
        for row, channel in enumerate(channels, start=1):
            # Channel name
            table.cell(row, 0).text = channel
            table.cell(row, 0).text_frame.paragraphs[0].font.size = Pt(11)

            # RMS value
            rms = rms_values.get(channel, 0.0)
            table.cell(row, 1).text = f"{rms:.4f}"
            table.cell(row, 1).text_frame.paragraphs[0].font.size = Pt(11)
            table.cell(row, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            # Additional metrics
            if additional_metrics:
                for col, metric_name in enumerate(additional_metrics.keys(), start=2):
                    value = additional_metrics[metric_name].get(channel, 0.0)
                    table.cell(row, col).text = f"{value:.4f}"
                    table.cell(row, col).text_frame.paragraphs[0].font.size = Pt(11)
                    table.cell(row, col).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        self._slide_count += 1

    def add_comparison_plot(
        self,
        image_bytes: bytes,
        title: str = "PSD Comparison",
        description: str = ""
    ) -> None:
        """
        Add a slide with a PSD comparison plot.

        Parameters
        ----------
        image_bytes : bytes
            PNG image data of the comparison plot.
        title : str, optional
            Slide title.
        description : str, optional
            Description text to display below the plot.
        """
        slide_layout = self.presentation.slide_layouts[6]  # Blank layout
        slide = self.presentation.slides.add_slide(slide_layout)

        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        self._style_slide_title(title_para)

        # Add the plot image
        self._add_picture_fit(slide, image_bytes, Inches(0.5), Inches(1.0), Inches(12.333), Inches(5.5))

        # Add description if provided
        if description:
            desc_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(6.6), Inches(12.333), Inches(0.5)
            )
            desc_frame = desc_box.text_frame
            desc_para = desc_frame.paragraphs[0]
            desc_para.text = description
            desc_para.font.size = Pt(12)
            desc_para.font.color.rgb = RGBColor(0x4a, 0x55, 0x68)
            desc_para.alignment = PP_ALIGN.CENTER

        self._slide_count += 1

    def add_text_slide(
        self,
        title: str,
        content: str,
        bullet_points: Optional[List[str]] = None
    ) -> None:
        """
        Add a text-only slide with optional bullet points.

        Parameters
        ----------
        title : str
            Slide title.
        content : str
            Main text content.
        bullet_points : list of str, optional
            List of bullet points to add.
        """
        slide_layout = self.presentation.slide_layouts[6]  # Blank layout
        slide = self.presentation.slides.add_slide(slide_layout)

        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        self._style_slide_title(title_para)

        # Add content
        content_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.2), Inches(12.333), Inches(5.5)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True

        # Main content
        para = content_frame.paragraphs[0]
        para.text = content
        para.font.size = Pt(14)
        para.font.color.rgb = RGBColor(0x4a, 0x55, 0x68)

        # Bullet points
        if bullet_points:
            for point in bullet_points:
                para = content_frame.add_paragraph()
                para.text = f"• {point}"
                para.font.size = Pt(12)
                para.font.color.rgb = RGBColor(0x4a, 0x55, 0x68)
                para.level = 0

        self._slide_count += 1

    def save(self, output_path: str) -> str:
        """
        Save the report to a PowerPoint file.

        Parameters
        ----------
        output_path : str
            Path where the PowerPoint file will be saved.

        Returns
        -------
        str
            The path where the file was saved.
        """
        # Ensure .pptx extension
        if not output_path.lower().endswith('.pptx'):
            output_path += '.pptx'

        self.presentation.save(output_path)
        return output_path

    def save_to_bytes(self) -> bytes:
        """
        Save the report to bytes (for in-memory handling).

        Returns
        -------
        bytes
            The PowerPoint file as bytes.
        """
        output = io.BytesIO()
        self.presentation.save(output)
        output.seek(0)
        return output.read()

    @property
    def slide_count(self) -> int:
        """Return the number of slides in the presentation."""
        return self._slide_count


def export_plot_to_image(plot_widget, width: int = 1200, height: int = 800) -> bytes:
    """
    Export a PyQtGraph plot widget to PNG image bytes.

    Parameters
    ----------
    plot_widget : pg.PlotWidget
        The PyQtGraph plot widget to export.
    width : int, optional
        Image width in pixels. Default is 1200.
    height : int, optional
        Image height in pixels. Default is 800.

    Returns
    -------
    bytes
        PNG image data.
    """
    try:
        from pyqtgraph.exporters import ImageExporter

        # Create exporter
        exporter = ImageExporter(plot_widget.plotItem)
        exporter.parameters()['width'] = width
        exporter.parameters()['height'] = height

        # Export to temporary file (ImageExporter.export() requires a filename)
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp_file:
            tmp_path = tmp_file.name

        try:
            exporter.export(tmp_path)

            # Read the bytes from the temporary file
            with open(tmp_path, 'rb') as f:
                image_bytes = f.read()

            return image_bytes

        finally:
            # Clean up the temporary file
            try:
                Path(tmp_path).unlink()
            except Exception:
                pass  # Ignore errors during cleanup

    except Exception as e:
        raise RuntimeError(f"Failed to export plot: {e}")
